from pptx import Presentation
from io import BytesIO
import xml.etree.ElementTree as ET
from xml.dom import minidom
import re


def extract_first_image_from_slide(pptx_file):
    """Extract the first image from the first slide of a PPTX file."""
    try:
        # Load the PowerPoint presentation from the uploaded file
        # pptx_file is expected to be a file-like object, like BytesIO
        ppt = Presentation(pptx_file)

        # Extract the first slide
        slide = ppt.slides[0]  # Change index if you want a specific slide

        # Iterate through shapes and find the first image
        for shape in slide.shapes:
            # Check if the shape contains an image (shape_type 13 indicates picture)
            if shape.shape_type == 13:
                image = shape.image

                # Convert the image blob to a BytesIO object
                image_stream = BytesIO(image.blob)

                return image_stream  # Return the image binary data

        raise Exception("No image found on the first slide.")

    except Exception as e:
        raise Exception(f"Error extracting image from PPTX: {e}")


def create_sample_data_with_header(input_file_path, data, output_file_path):
    # Read the input file
    with open(input_file_path, 'r', encoding='utf-8') as file:
        content = file.read()

    # Extract the DOCTYPE declaration (header)
    start = content.find('<?xml')
    end = content.find(']>') + 2
    header = content[start:end]
    body = content[end:].strip()

    # Locate the <v:sampleDataSets> element
    sample_data_sets_start = body.find('<v:sampleDataSets')
    sample_data_sets_end = body.find(
        '</v:sampleDataSets>') + len('</v:sampleDataSets>')
    sample_data_sets_content = body[sample_data_sets_start:sample_data_sets_end]

    # Clear the content inside <v:sampleDataSets> before adding new data
    # Extract the opening tag
    opening_tag_end = sample_data_sets_content.find('>') + 1
    opening_tag = sample_data_sets_content[:opening_tag_end]

    # Create a new <v:sampleDataSets> element with new data
    modified_sample_data_sets = f"""{opening_tag}
    <v:sampleDataSet dataSetName="1">
"""
    for key, value in data.items():
        modified_sample_data_sets += f"""
        <{key}>
            <p>{value}</p>
        </{key}>
"""
    modified_sample_data_sets += """
    </v:sampleDataSet>
</v:sampleDataSets>
"""

    # Replace the old <v:sampleDataSets> with the new one
    body = body.replace(sample_data_sets_content, modified_sample_data_sets)

    # Write the updated content to the output file
    with open(output_file_path, 'w', encoding='utf-8') as file:
        file.write(header)
        file.write('\n')
        file.write(body)
