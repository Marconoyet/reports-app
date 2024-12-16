
from pptx import Presentation
from io import BytesIO
import xml.etree.ElementTree as ET
from db.reports import get_report_db
import base64
from spire.pdf.common import *
from spire.pdf import *
from io import BytesIO
from db.custom_exceptions import DatabaseError
from db.reports import (
    add_report_db,
    add_report_xml_db,
    get_report_db,
    delete_report_db,
    update_report_db,
    get_file_of_report
)
from db.files import upload_file_db, get_report_with_template_data
from sqlalchemy.exc import SQLAlchemyError
from services.services_utiliy import extract_first_image_from_slide, create_sample_data_with_header
from datetime import datetime
import threading
import subprocess
import tempfile
import os
import rarfile
from zipfile import ZipFile
from bs4 import BeautifulSoup
from flask import current_app
import mimetypes


def process_rar_file(rar_file):
    try:
        with rarfile.RarFile(rar_file) as rf:
            html_file = None
            images = {}
            for file in rf.namelist():
                if file.endswith(".html"):
                    html_file = rf.read(file).decode("utf-8")
                elif file.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.svg')):
                    images[file] = base64.b64encode(
                        rf.read(file)).decode('utf-8')

            if not html_file:
                raise Exception("No HTML file found in RAR archive")

            return embed_images_in_html(html_file, images)
    except rarfile.Error as e:
        raise Exception(f"Error processing RAR file: {e}")


def process_zip_file(zip_file):
    try:
        # Unzip the file into memory
        with ZipFile(zip_file) as z:
            # Extract the HTML file and images
            html_file = None
            images = {}
            for filename in z.namelist():
                if filename.endswith(".html"):
                    html_file = z.read(filename).decode('utf-8')
                elif filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.svg')):
                    images[filename] = base64.b64encode(
                        z.read(filename)).decode('utf-8')

            if not html_file:
                raise Exception("No HTML file found in ZIP")

            # Embed images into HTML
            return embed_images_in_html(html_file, images)

    except Exception as e:
        raise Exception(f"Error processing ZIP file: {e}")


def embed_images_in_html(html_content, images):
    soup = BeautifulSoup(html_content, "html.parser")

    for img_tag in soup.find_all("img"):
        img_src = img_tag.get("src")
        if img_src in images:
            # Detect MIME type based on the file extension
            mime_type, _ = mimetypes.guess_type(img_src)
            if not mime_type:
                mime_type = "application/octet-stream"  # Default MIME type
            # Replace src with Base64-encoded data
            img_tag["src"] = f"data:{mime_type};base64,{images[img_src]}"
        else:
            print(f"Warning: Image {img_src} not found in ZIP.")

    return str(soup)


def allowed_file(filename):
    allowed_extensions = {'pptx'}
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in allowed_extensions


def add_report_xml(report_data):
    """Business logic for retrieving a report."""
    try:
        return add_report_xml_db(report_data)
    except DatabaseError as e:
        raise Exception(f"Service Error - Could not retrieve report: {e}")


def add_report(report_data):
    """Business logic for adding a report and saving the first image from the first slide in the template."""
    try:
        file = report_data.get('file')

        print("File saved as report.pptx")
        if file and allowed_file(file.filename):
            # Step 1: Extract the first image from the first slide
            file_binary_data = file.read()
            file_stream = BytesIO(file_binary_data).getvalue()
            first_image_stream = extract_first_image_from_slide(file)

            # Step 2: Convert the image to binary data
            image_bytes = first_image_stream.getvalue()

            # Step 3: Insert report metadata into MySQL
            report_metadata = {
                "template_name": report_data.get('reportName'),
                "template_description": report_data.get('reportDescription'),
                "template_image": image_bytes,
                "user_id": int(report_data.get('userId')),
                "folder_id": int(report_data.get('folderId')),
                "center_id": int(report_data.get('centerId')),
                "template_file": file_stream,
                "created_time": datetime.now()
            }

            # Step 4: Insert the report into the MySQL database
            report_id = add_report_db(report_metadata)

            return {"status": "success", "report_id": report_id}

        else:
            raise Exception("Invalid file type or no file uploaded.")

    except SQLAlchemyError as e:
        raise Exception(f"Database operation failed: {e}")

    except Exception as e:
        raise Exception(f"Service Error - Could not add report: {e}")


def get_report(report_id):
    """Business logic for retrieving a report."""
    try:
        return get_report_db(report_id)
    except DatabaseError as e:
        raise Exception(f"Service Error - Could not retrieve report: {e}")


def delete_report(report_id):
    """Business logic for deleting a report."""
    try:
        return delete_report_db(report_id)
    except DatabaseError as e:
        raise Exception(f"Service Error - Could not delete report: {e}")


def update_report(report_id, updated_data):
    """Business logic for deleting a report."""
    try:
        return update_report_db(report_id, updated_data)
    except DatabaseError as e:
        raise Exception(f"Service Error - Could not delete report: {e}")


def generate_xml_report(replacements, report_id, report_name):

    try:
        # Step 1: Retrieve the report and file data from the database
        report = get_file_of_report(report_id)
        file_data = report.template_file  # XML template file from the database

        if not file_data:
            raise Exception("No template file found for the specified report.")

        # Step 2: Create temporary input and output files for processing
        with tempfile.NamedTemporaryFile(suffix=".xml", delete=False) as temp_input_file:
            temp_input_file.write(file_data)
            temp_input_path = temp_input_file.name

        with tempfile.NamedTemporaryFile(suffix=".xml", delete=False) as temp_output_file:
            temp_output_path = temp_output_file.name

        # Step 3: Apply replacements to generate the updated XML content
        create_sample_data_with_header(
            temp_input_path, replacements, temp_output_path)

        # Step 4: Read the modified file and return it as a BytesIO stream
        with open(temp_output_path, 'rb') as modified_file:
            modified_file_stream = BytesIO(modified_file.read())
        app = current_app._get_current_object()
        threading.Thread(target=upload_pptx_in_background,
                         args=(modified_file_stream, report, app, report_name, True)).start()
        return modified_file_stream
    except Exception as e:
        raise Exception(f"Error generating XML report: {e}")

    finally:
        # Cleanup temporary files
        try:
            if temp_input_path:
                os.remove(temp_input_path)
            if temp_output_path:
                os.remove(temp_output_path)
        except Exception as cleanup_error:
            print(f"Cleanup error: {cleanup_error}")


def generate_pptx_report(replacements, report_id, report_name):

    # Generate PPTX report
    report = get_file_of_report(report_id)
    file_data = report.template_file
    presentation = Presentation(BytesIO(file_data))

    # Replace the fields in the PPTX
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key, value in replacements.items():
                        if key in run.text:
                            run.text = run.text.replace(key, value)

    # Save the modified presentation to a BytesIO stream
    pptx_stream = BytesIO()
    presentation.save(pptx_stream)
    pptx_stream.seek(0)

    # Save the PPTX stream to a temporary file for PDF conversion
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_pptx_file:
        temp_pptx_file.write(pptx_stream.getvalue())
        temp_pptx_path = temp_pptx_file.name

    pdf_stream = None
    pdf_error = None  # To store error messages

    # Convert the PPTX to PDF
    with tempfile.TemporaryDirectory() as temp_output_dir:
        try:
            env = os.environ.copy()
            env["PATH"] += os.pathsep + "/usr/bin"
            libreoffice_path = '/usr/bin/libreoffice'  # Update this path if necessary
            command = [
                libreoffice_path, '--headless', '--convert-to', 'pdf', temp_pptx_path, '--outdir', temp_output_dir
            ]
            subprocess.run(command, check=True, env=env)
            pdf_path = os.path.join(temp_output_dir, os.path.splitext(
                os.path.basename(temp_pptx_path))[0] + ".pdf")

            # Read the generated PDF into a BytesIO stream
            pdf_stream = BytesIO()
            with open(pdf_path, 'rb') as pdf_file:
                pdf_stream.write(pdf_file.read())
            pdf_stream.seek(0)
        except subprocess.CalledProcessError as e:
            pdf_error = f"Error during conversion: {e}"
        except FileNotFoundError as e:
            pdf_error = f"File not found error: {e}"
        finally:
            # Cleanup temporary PPTX file
            try:
                os.remove(temp_pptx_path)
            except OSError as cleanup_error:
                pdf_error = f"Error cleaning up temporary file: {cleanup_error}"

    # Start a background thread to handle the PPTX file upload process
    app = current_app._get_current_object()
    threading.Thread(target=upload_pptx_in_background,
                     args=(pptx_stream, report, app, report_name)).start()

    # Return both PPTX and PDF streams, and the error if present
    return pptx_stream, pdf_stream, pdf_error


def upload_pptx_in_background(output_stream, report, app, report_name, xml=False):
    # Prepare the report for upload
    # Read the stream into binary for upload
    report.template_file = output_stream.getvalue()
    neededData = prepareReportForUpload(report, report_name, xml)
    with app.app_context():
        try:
            upload_response = upload_file_db(neededData)
            print(f"File upload successful: {upload_response}")
        except Exception as e:
            print(f"Error uploading file in background: {e}")


def prepareReportForUpload(report, report_name, xml=False):
    try:
        # Extract the file from the report
        image_bytes = None
        report_data = report.to_dict()
        file = report.template_file
        if not file:
            print("The report has no associated template file.")
            raise ValueError("The report has no associated template file.")
        if isinstance(file, bytes):
            file_binary_data = file
            # Wrap bytes in BytesIO for further operations
            file_io = BytesIO(file_binary_data)
        else:
            raise ValueError("Unexpected file format")
        if (xml != True):
            first_image_stream = extract_first_image_from_slide(file_io)
            image_bytes = first_image_stream.getvalue()
    except AttributeError as e:
        print(
            f"Error accessing the report file or file-related attributes: {e}")
        raise Exception(
            f"Error accessing the report file or file-related attributes: {e}")
    except Exception as e:
        print(f"Error processing the template file: {e}")
        raise Exception(f"Error processing the template file: {e}")

    try:
        # Prepare the report metadata
        report_metadata = {
            "report_name": report_name,  # Report name
            "template_id": report_data.get("id"),  # Template ID (foreign key)
            "report_file": file_binary_data,                            # Binary file data
            "report_image": image_bytes,
            "center_id": report_data.get("center_id"),
            "user_id": report_data.get("user_id"),
            "created_time": datetime.now()                  # Current timestamp
        }

        return report_metadata

    except KeyError as e:
        print(e)
        raise Exception(f"Missing required report data: {e}")
    except Exception as e:
        print(e)
        raise Exception(
            f"An unexpected error occurred while preparing report metadata: {e}")


def handle_extract_xml_fields(report_id, viewer=False):
    try:
        # Step 1: Retrieve the report file from the database
        print(viewer)
        if viewer:
            report_data = get_report_with_template_data(report_id)

        else:
            report = get_file_of_report(report_id)
            file_data = report.template_file  # Binary data of the file
            # Decode the binary HTML data to a UTF-8 string (if it is an HTML file)
            html_file_content = file_data.decode(
                'utf-8') if file_data else None
            # Step 3: Prepare report data
            report_data = report.to_dict()
            if report.template_image:
                encoded_image = base64.b64encode(
                    report.template_image).decode('utf-8')
                report_data['template_image'] = f"data:image/png;base64,{encoded_image}"

            # Include the HTML content directly
            report_data["template_file"] = html_file_content

            # Step 4: Combine fields and report data
        combined_data = {
            "report": report_data,  # Report details
        }
        return combined_data
    except Exception as e:
        raise Exception(
            f"Service Error - Could not extract fields and report data: {e}")


def extract_varnames_from_svg_file(file_stream):
    """
    Extract variable names (varName attributes) from an SVG file.
    Args:
        file_stream: A file-like object (BytesIO) containing the SVG data.
    Returns:
        A list of variable names found in the SVG file.
    """
    try:
        # Parse the SVG content
        tree = ET.parse(file_stream)
        root = tree.getroot()

        # Define the namespace for Adobe Variables
        namespace = {'ns_vars': "http://ns.adobe.com/Variables/1.0/"}

        # Find all variable elements
        variables = root.findall(
            './/ns_vars:variables/ns_vars:variable', namespaces=namespace)

        # Extract the 'varName' attribute values
        varnames = [var.get('varName')
                    for var in variables if var.get('varName')]

        return varnames
    except ET.ParseError as e:
        raise Exception(f"Failed to parse SVG file: {e}")


def get_report_and_extract_fields(report_id):
    """Service layer to extract fields from PPTX or PDF files."""
    try:
        fields = extract_pptx_fields(report_id)
        return fields
    except Exception as e:
        raise Exception(f"Service Error - Could not extract fields: {e}")


def extract_pptx_fields(report_id):
    """Extract fields from a PPTX file."""
    report = get_file_of_report(report_id)
    file_data = report.template_file
    presentation = Presentation(BytesIO(file_data))
    fields = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "#" in run.text:  # Assume fields start with '#'
                            fields.append(run.text)
    report_data = report.to_dict()
    if report.template_image:
        encoded_image = base64.b64encode(report.template_image).decode('utf-8')
        report_data['template_image'] = f"data:image/png;base64,{encoded_image}"
    combined_data = {
        "report": report_data,            # Folder details
        "fields": fields  # Templates associated with the folder
    }
    return combined_data
